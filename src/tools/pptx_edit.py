#!/usr/bin/env python3
"""Server-side PPTX editor invoked by the microclaw `pptx_edit` tool.

Reads a JSON spec on stdin, performs structured edits using python-pptx,
writes a JSON result to stdout. Stays small and dependency-light on
purpose — the only runtime dep is `python-pptx`, which is already in
the microclaw image.

JSON spec shape:
  {
    "source_path":  "/abs/path/in.pptx",
    "output_path":  "/abs/path/out.pptx",   # omit to inspect only
    "operations":   [ {op, ...}, ... ]      # optional; if absent, just inspect
  }

Result shape:
  {
    "ok": bool,
    "mode": "inspect" | "apply",
    "source": {...slide summary...},
    "operations_applied": int,
    "errors": [str, ...],
    "output_path": str | null
  }
"""
from __future__ import annotations

import json
import sys
import traceback
from copy import deepcopy

try:
    from pptx import Presentation
    from pptx.util import Pt, Emu
    from pptx.dml.color import RGBColor
except ImportError as e:
    print(json.dumps({
        "ok": False,
        "errors": [f"python-pptx not installed: {e}"],
    }))
    sys.exit(2)


# ---------- inspection ----------

def inspect_presentation(prs):
    """Return a JSON-friendly summary of a Presentation."""
    slides = []
    for idx, slide in enumerate(prs.slides):
        layout_name = slide.slide_layout.name if slide.slide_layout else "(unknown)"
        shapes_summary = []
        for s_idx, shape in enumerate(slide.shapes):
            sd = {
                "shape_index": s_idx,
                "name": shape.name,
                "shape_type": str(shape.shape_type),
                "is_placeholder": bool(getattr(shape, "is_placeholder", False)),
                "has_text_frame": bool(getattr(shape, "has_text_frame", False)),
            }
            if getattr(shape, "is_placeholder", False) and shape.placeholder_format is not None:
                sd["placeholder_idx"] = shape.placeholder_format.idx
                sd["placeholder_type"] = str(shape.placeholder_format.type)
            if getattr(shape, "has_text_frame", False):
                paras = []
                for p_idx, para in enumerate(shape.text_frame.paragraphs):
                    paras.append({
                        "paragraph_index": p_idx,
                        "level": para.level,
                        "text": para.text,
                    })
                sd["paragraphs"] = paras
            shapes_summary.append(sd)
        title = ""
        try:
            if slide.shapes.title is not None:
                title = slide.shapes.title.text or ""
        except Exception:
            pass
        slides.append({
            "slide_index": idx,
            "layout_name": layout_name,
            "title": title,
            "shape_count": len(slide.shapes),
            "shapes": shapes_summary,
        })
    return {
        "slide_count": len(prs.slides),
        "slide_width_emu": prs.slide_width,
        "slide_height_emu": prs.slide_height,
        "slides": slides,
    }


# ---------- helpers ----------

def _get_slide(prs, slide_idx):
    if slide_idx < 0 or slide_idx >= len(prs.slides):
        raise IndexError(f"slide {slide_idx} out of range (have {len(prs.slides)} slides)")
    return prs.slides[slide_idx]


def _replace_text_in_text_frame(tf, find, replace, case_sensitive):
    """Replace text run-by-run so formatting is preserved where the
    match is contained inside a single run. Multi-run matches fall
    back to whole-paragraph replacement (loses run-level formatting
    inside that paragraph)."""
    n = 0
    for para in tf.paragraphs:
        for run in para.runs:
            haystack = run.text or ""
            if case_sensitive:
                if find in haystack:
                    run.text = haystack.replace(find, replace)
                    n += haystack.count(find)
            else:
                if find.lower() in haystack.lower():
                    # Case-insensitive but preserve original casing for non-matches.
                    new = []
                    i = 0
                    f_len = len(find)
                    lower = haystack.lower()
                    f_lower = find.lower()
                    while i < len(haystack):
                        if lower[i:i+f_len] == f_lower:
                            new.append(replace)
                            i += f_len
                            n += 1
                        else:
                            new.append(haystack[i])
                            i += 1
                    run.text = "".join(new)
        # Multi-run match fallback: if the joined paragraph text contains the
        # needle but no single run did, replace at the paragraph level.
        joined = para.text or ""
        target = joined if case_sensitive else joined.lower()
        needle = find if case_sensitive else find.lower()
        if needle in target and not any(
            (find in (r.text or "")) if case_sensitive else (find.lower() in (r.text or "").lower())
            for r in para.runs
        ):
            new_text = joined.replace(find, replace) if case_sensitive else \
                       joined.lower().replace(needle, replace)
            # Wipe runs and set text on first run.
            if para.runs:
                para.runs[0].text = new_text
                for r in para.runs[1:]:
                    r.text = ""
            n += joined.count(find) if case_sensitive else target.count(needle)
    return n


# ---------- operations ----------

def op_replace_text(prs, op):
    find = op["find"]
    replace = op.get("replace", "")
    case_sensitive = op.get("case_sensitive", False)
    slide_filter = op.get("slide")  # optional int; None = all
    n = 0
    for s_idx, slide in enumerate(prs.slides):
        if slide_filter is not None and s_idx != slide_filter:
            continue
        for shape in slide.shapes:
            if getattr(shape, "has_text_frame", False):
                n += _replace_text_in_text_frame(shape.text_frame, find, replace, case_sensitive)
    return {"replacements": n}


def op_set_slide_title(prs, op):
    slide = _get_slide(prs, op["slide"])
    title_shape = slide.shapes.title
    if title_shape is None:
        raise ValueError(f"slide {op['slide']} has no title placeholder")
    title_shape.text = op["title"]
    return {"set": True}


def op_set_paragraph_text(prs, op):
    slide = _get_slide(prs, op["slide"])
    shape_idx = op["shape"]
    para_idx = op.get("paragraph", 0)
    new_text = op["text"]
    if shape_idx < 0 or shape_idx >= len(slide.shapes):
        raise IndexError(f"shape {shape_idx} out of range on slide {op['slide']}")
    shape = list(slide.shapes)[shape_idx]
    if not getattr(shape, "has_text_frame", False):
        raise ValueError("shape has no text frame")
    paras = list(shape.text_frame.paragraphs)
    if para_idx < 0 or para_idx >= len(paras):
        raise IndexError(f"paragraph {para_idx} out of range")
    paras[para_idx].text = new_text
    return {"set": True}


def op_add_bullet(prs, op):
    slide = _get_slide(prs, op["slide"])
    shape_idx = op["shape"]
    text = op["text"]
    level = op.get("level", 0)
    if shape_idx < 0 or shape_idx >= len(slide.shapes):
        raise IndexError(f"shape {shape_idx} out of range on slide {op['slide']}")
    shape = list(slide.shapes)[shape_idx]
    if not getattr(shape, "has_text_frame", False):
        raise ValueError("shape has no text frame")
    p = shape.text_frame.add_paragraph()
    p.text = text
    p.level = level
    return {"added": True}


def op_delete_slide(prs, op):
    slide_idx = op["slide"]
    if slide_idx < 0 or slide_idx >= len(prs.slides):
        raise IndexError(f"slide {slide_idx} out of range")
    # python-pptx has no delete_slide, so we manipulate the XML directly.
    xml_slides = prs.slides._sldIdLst  # noqa: SLF001
    slides = list(xml_slides)
    xml_slides.remove(slides[slide_idx])
    return {"deleted": slide_idx}


def op_add_slide(prs, op):
    layout_idx = op.get("layout", 1)
    after_slide = op.get("after_slide")  # None = append; -1 = beginning
    if layout_idx < 0 or layout_idx >= len(prs.slide_layouts):
        raise IndexError(f"layout {layout_idx} out of range (have {len(prs.slide_layouts)})")
    layout = prs.slide_layouts[layout_idx]
    new_slide = prs.slides.add_slide(layout)
    if op.get("title"):
        if new_slide.shapes.title is not None:
            new_slide.shapes.title.text = op["title"]
    body = op.get("body") or []
    if body:
        # Find a body placeholder.
        body_shape = None
        for shape in new_slide.placeholders:
            if shape.placeholder_format.idx != 0 and getattr(shape, "has_text_frame", False):
                body_shape = shape
                break
        if body_shape is not None:
            tf = body_shape.text_frame
            tf.text = body[0] if body else ""
            for line in body[1:]:
                p = tf.add_paragraph()
                p.text = line
    # Reposition if after_slide given.
    if after_slide is not None:
        xml_slides = prs.slides._sldIdLst  # noqa: SLF001
        slides_xml = list(xml_slides)
        last = slides_xml[-1]  # the new one (we just added)
        xml_slides.remove(last)
        target_pos = after_slide + 1
        target_pos = max(0, min(target_pos, len(slides_xml)))
        xml_slides.insert(target_pos, last)
    return {"added": True, "slide_index": len(prs.slides) - 1}


def op_reorder_slides(prs, op):
    new_order = op["order"]
    if sorted(new_order) != list(range(len(prs.slides))):
        raise ValueError(
            f"order must be a permutation of 0..{len(prs.slides) - 1}; got {new_order}"
        )
    xml_slides = prs.slides._sldIdLst  # noqa: SLF001
    slides_xml = list(xml_slides)
    for s in slides_xml:
        xml_slides.remove(s)
    for new_idx in new_order:
        xml_slides.append(slides_xml[new_idx])
    return {"reordered": new_order}


def op_set_font_size(prs, op):
    slide = _get_slide(prs, op["slide"])
    shape_idx = op["shape"]
    size_pt = op["size"]
    para_idx = op.get("paragraph")  # optional; None = all paragraphs
    shape = list(slide.shapes)[shape_idx]
    if not getattr(shape, "has_text_frame", False):
        raise ValueError("shape has no text frame")
    paras = list(shape.text_frame.paragraphs)
    targets = paras if para_idx is None else [paras[para_idx]]
    n = 0
    for p in targets:
        for r in p.runs:
            r.font.size = Pt(size_pt)
            n += 1
    return {"runs_set": n}


def op_set_font_color(prs, op):
    slide = _get_slide(prs, op["slide"])
    shape_idx = op["shape"]
    rgb_hex = op["rgb"].lstrip("#")
    if len(rgb_hex) != 6:
        raise ValueError(f"rgb must be 6 hex chars; got {op['rgb']!r}")
    rgb = RGBColor.from_string(rgb_hex)
    para_idx = op.get("paragraph")
    shape = list(slide.shapes)[shape_idx]
    if not getattr(shape, "has_text_frame", False):
        raise ValueError("shape has no text frame")
    paras = list(shape.text_frame.paragraphs)
    targets = paras if para_idx is None else [paras[para_idx]]
    n = 0
    for p in targets:
        for r in p.runs:
            r.font.color.rgb = rgb
            n += 1
    return {"runs_set": n}


OP_REGISTRY = {
    "replace_text":       op_replace_text,
    "set_slide_title":    op_set_slide_title,
    "set_paragraph_text": op_set_paragraph_text,
    "add_bullet":         op_add_bullet,
    "delete_slide":       op_delete_slide,
    "add_slide":          op_add_slide,
    "reorder_slides":     op_reorder_slides,
    "set_font_size":      op_set_font_size,
    "set_font_color":     op_set_font_color,
}


def main():
    try:
        spec = json.load(sys.stdin)
    except Exception as e:
        print(json.dumps({"ok": False, "errors": [f"invalid JSON spec: {e}"]}))
        sys.exit(1)

    source_path = spec.get("source_path")
    output_path = spec.get("output_path")
    operations = spec.get("operations") or []

    if not source_path:
        print(json.dumps({"ok": False, "errors": ["source_path required"]}))
        sys.exit(1)

    try:
        prs = Presentation(source_path)
    except Exception as e:
        print(json.dumps({
            "ok": False,
            "errors": [f"failed to open {source_path}: {e}"],
        }))
        sys.exit(1)

    source_summary = inspect_presentation(prs)

    # Inspect-only mode: no operations + no output_path → return summary.
    if not operations and not output_path:
        print(json.dumps({
            "ok": True,
            "mode": "inspect",
            "source": source_summary,
            "operations_applied": 0,
            "errors": [],
            "output_path": None,
        }))
        return

    if not output_path:
        print(json.dumps({
            "ok": False,
            "errors": ["output_path required when operations are given"],
        }))
        sys.exit(1)

    # Apply mode.
    errors = []
    applied = 0
    op_results = []
    for i, op in enumerate(operations):
        op_name = op.get("op")
        if op_name not in OP_REGISTRY:
            errors.append(f"op[{i}]: unknown op {op_name!r}")
            continue
        try:
            result = OP_REGISTRY[op_name](prs, op)
            op_results.append({"index": i, "op": op_name, "result": result})
            applied += 1
        except Exception as e:
            errors.append(f"op[{i}] ({op_name}): {e.__class__.__name__}: {e}")
            op_results.append({"index": i, "op": op_name, "error": str(e)})

    try:
        prs.save(output_path)
    except Exception as e:
        print(json.dumps({
            "ok": False,
            "errors": errors + [f"failed to save {output_path}: {e}"],
            "operations_applied": applied,
            "op_results": op_results,
        }))
        sys.exit(1)

    print(json.dumps({
        "ok": True,
        "mode": "apply",
        "source": {
            "slide_count": source_summary["slide_count"],
        },
        "output_path": output_path,
        "operations_applied": applied,
        "op_results": op_results,
        "errors": errors,
    }))


if __name__ == "__main__":
    try:
        main()
    except Exception as e:
        print(json.dumps({
            "ok": False,
            "errors": [f"unexpected error: {e}", traceback.format_exc()],
        }))
        sys.exit(1)
